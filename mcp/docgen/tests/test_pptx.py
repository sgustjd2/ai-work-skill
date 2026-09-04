import re
import zipfile

import pytest
from pptx import Presentation

from docgen import parse, pptx_render

pytestmark = pytest.mark.deterministic


@pytest.fixture
def rendered(deck_md, theme, tmp_path):
    parsed = parse.parse_deck(deck_md)
    out = tmp_path / "deck.pptx"
    res = pptx_render.render_pptx(parsed, str(out), theme, base_dir=str(tmp_path))
    return out, res


def _slide_xmls(path):
    with zipfile.ZipFile(path) as z:
        return [
            z.read(n).decode("utf-8")
            for n in z.namelist()
            if re.match(r"ppt/slides/slide\d+\.xml", n)
        ]


def test_slide_count_cover_agenda(rendered):
    out, res = rendered
    # 6 content + cover + agenda = 8
    assert res["slide_count"] == 8
    assert len(Presentation(str(out)).slides) == 8


def test_layouts_recorded(rendered):
    _, res = rendered
    layouts = {s["title"]: s["layout"] for s in res["slides"]}
    assert layouts["목표 아키텍처"] == "diagram"
    assert layouts["비용 비교"] == "chart"


def test_headline_text_present(rendered):
    out, _ = rendered
    xmls = " ".join(_slide_xmls(out))
    assert "게이트웨이 한 곳에서" in xmls


def test_eastasia_font(rendered):
    out, _ = rendered
    assert any("a:ea" in x for x in _slide_xmls(out))


def test_arrow_tailend(rendered):
    out, _ = rendered
    assert any("tailEnd" in x for x in _slide_xmls(out))


def test_native_chart(rendered):
    out, _ = rendered
    with zipfile.ZipFile(str(out)) as z:
        assert any(n.startswith("ppt/charts/chart") for n in z.namelist())


def test_notes(rendered):
    out, _ = rendered
    with zipfile.ZipFile(str(out)) as z:
        assert any("notesSlide" in n for n in z.namelist())


def test_footer_page_number(rendered):
    out, _ = rendered
    xmls = " ".join(_slide_xmls(out))
    assert re.search(r"\d+ / \d+", xmls)


def test_template_mode(theme, tmp_path):
    tpl = Presentation()
    tpl.slides.add_slide(tpl.slide_layouts[0])  # 기존 슬라이드 1장
    tpl_path = tmp_path / "company.pptx"
    tpl.save(str(tpl_path))

    parsed = parse.parse_deck("---\ntitle: t\n---\n# s\n## 헤드 메시지\n- 항목\n")
    out = tmp_path / "from_tpl.pptx"
    res = pptx_render.render_pptx(
        parsed, str(out), theme, base_dir=str(tmp_path), template_pptx=str(tpl_path)
    )
    # 템플릿 기존 슬라이드는 제거되고 cover + 1 content = 2
    assert res["slide_count"] == 2
