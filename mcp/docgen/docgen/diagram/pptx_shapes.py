"""배치 → pptx 네이티브 도형(편집 가능). render_pptx 의 diagram 레이아웃이 쓴다."""

from __future__ import annotations

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from .. import theme as T
from .layout import KIND_STYLE, build_layout


def _rgb(theme, role):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(T.rgb_hex(theme, role))


def _rgb_hex(theme, kind):
    from pptx.dml.color import RGBColor

    role = KIND_STYLE[kind]["fill"]
    hexv = T.tint(theme, "accent", 0.9) if role == "llm_tint" else T.color(theme, role)
    return RGBColor.from_string(hexv.lstrip("#").upper())


def draw_diagram(slide, spec, theme, rect):
    """rect = (left_in, top_in, width_in, height_in). 비율 유지해 맞춘다."""
    lay = build_layout(spec)
    left, top, rw, rh = rect
    lw, lh = lay["width"], lay["height"]
    scale = min(rw / lw, rh / lh) if lw and lh else 1.0
    offx = left + (rw - lw * scale) / 2
    offy = top + (rh - lh * scale) / 2

    def X(v):
        return Inches(offx + v * scale)

    def Y(v):
        return Inches(offy + v * scale)

    def S(v):
        return Inches(v * scale)

    ink = _rgb(theme, "ink")

    for g in lay["groups"]:
        if not g["label"]:
            continue
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, X(g["x"]), Y(g["y"]), S(g["w"]), S(g["h"])
        )
        box.name = f"group:{g['id']}"
        box.fill.background()
        hi = str(g.get("style", "")).lower() == "highlight"
        box.line.color.rgb = _rgb(theme, "accent" if hi else "line")
        box.line.width = Pt(1.25)
        if hi:
            _dash(box)
        box.shadow.inherit = False
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = g["label"]
        _run_font(r, theme, "ink_muted", 10)

    for e in lay["edges"]:
        kind = MSO_CONNECTOR.ELBOW if e.get("elbow") else MSO_CONNECTOR.STRAIGHT
        conn = slide.shapes.add_connector(kind, X(e["x1"]), Y(e["y1"]), X(e["x2"]), Y(e["y2"]))
        conn.line.color.rgb = ink
        conn.line.width = Pt(1.25)
        if str(e.get("style", "")).lower() == "dashed":
            _dash(conn)
        _arrow_end(conn)
        if e["label"]:
            mx, my = (e["x1"] + e["x2"]) / 2, (e["y1"] + e["y2"]) / 2
            tb = slide.shapes.add_textbox(
                X(mx) - Inches(0.5), Y(my) - Inches(0.12), Inches(1.0), Inches(0.24)
            )
            tb.fill.solid()
            tb.fill.fore_color.rgb = _rgb(theme, "white")
            tb.line.fill.background()
            tf = tb.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = e["label"]
            _run_font(r, theme, "ink_muted", 9)

    for n in lay["nodes"]:
        st = KIND_STYLE[n["kind"]]
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, X(n["x"]), Y(n["y"]), S(n["w"]), S(n["h"])
        )
        box.name = f"node:{n['id']}"
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb_hex(theme, n["kind"])
        box.line.color.rgb = _rgb(theme, st["stroke"])
        box.line.width = Pt(1.5)
        box.shadow.inherit = False
        if st["dashed"]:
            _dash(box)
        _adjust(box, 0.12)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Emu(36000)
        tf.margin_top = tf.margin_bottom = Emu(18000)
        for i, ln in enumerate(n["lines"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = ln
            _run_font(r, theme, st["text"], 12)

    return {"warnings": lay["warnings"]}


def _run_font(run, theme, role, size):
    run.font.size = Pt(size)
    run.font.color.rgb = _rgb(theme, role)
    run.font.name = T.font(theme, "office_latin")
    _set_ea(run, T.font(theme, "office_ko"))


def _set_ea(run, ko_name):
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", ko_name)
    rpr.set("lang", "ko-KR")


def _dash(shape):
    ln = shape.line._get_or_add_ln()
    d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
    ln.append(d)


def _adjust(shape, val):
    try:
        shape.adjustments[0] = val
    except (IndexError, ValueError):
        pass


def _arrow_end(conn):
    ln = conn.line._get_or_add_ln()
    end = ln.makeelement(qn("a:tailEnd"), {"type": "triangle"})
    ln.append(end)
