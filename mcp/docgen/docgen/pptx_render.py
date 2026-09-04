"""블록 모델(덱) → pptx (FR-11). 슬라이드마다 헤드 메시지 1개, 개조식, 테마·회사 템플릿만.

레이아웃: cover·agenda·section·message·two-col·diagram·table·timeline·chart·image·closing.
색·폰트·치수는 테마에서만. 한글은 run 마다 a:latin + a:ea 를 지정한다.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from . import theme as T
from .diagram import pptx_shapes

_BLANK = 6


def _rgb(theme, role):
    return RGBColor.from_string(T.rgb_hex(theme, role))


def _set_ea(run, ko_name):
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", ko_name)
    rpr.set("lang", "ko-KR")


class DeckRenderer:
    def __init__(self, theme, base_dir=None):
        self.theme = theme
        self.base_dir = Path(base_dir) if base_dir else Path.cwd()
        self.warnings: list[str] = []
        p = theme["pptx"]
        self.W = p["width_in"]
        self.H = p["height_in"]
        self.M = p.get("margin_in", 0.5)
        self.ko = T.font(theme, "office_ko")
        self.latin = T.font(theme, "office_latin")

    # ── 텍스트 헬퍼 ──
    def _text(
        self, tf, text, role, size, *, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP
    ):
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        self._font(r, role, size, bold)
        return p

    def _font(self, run, role, size, bold=False):
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(self.theme, role)
        run.font.name = self.latin
        _set_ea(run, self.ko)

    def _box(self, slide, left, top, w, h):
        return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))

    def _blank(self, prs):
        return prs.slides.add_slide(prs.slide_layouts[_BLANK])

    # ── 공통 프레임(제목·헤드 메시지·accent bar·바닥글) ──
    def _frame(self, slide, sl, index, total):
        p = self.theme["pptx"]
        content_top = self.M
        if sl.get("title"):
            tb = self._box(slide, self.M, self.M, self.W - 2 * self.M, 0.4)
            self._text(
                tb.text_frame, sl["title"], p.get("title_color", "ink_muted"), p.get("title_pt", 14)
            )
            content_top = self.M + 0.5
        if sl.get("headline"):
            hb = self._box(slide, self.M, content_top, self.W - 2 * self.M, 1.0)
            self._text(
                hb.text_frame,
                sl["headline"],
                p.get("headline_color", "primary"),
                p.get("headline_pt", 22),
                bold=True,
            )
            content_top += 1.0
            if p.get("accent_bar", True):
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(self.M),
                    Inches(content_top),
                    Inches(2.0),
                    Inches(0.04),
                )
                bar.fill.solid()
                bar.fill.fore_color.rgb = _rgb(self.theme, "primary")
                bar.line.fill.background()
                bar.shadow.inherit = False
                content_top += 0.2
        self._footer(slide, sl, index, total)
        return content_top, self.H - self.M - 0.35

    def _footer(self, slide, sl, index, total):
        p = self.theme["pptx"]
        y = self.H - self.M - 0.05
        fs = p.get("footer_pt", 9)
        left = self.theme.get("_footer_text", "")
        lb = self._box(slide, self.M, y - 0.1, self.W / 3, 0.3)
        self._text(lb.text_frame, left, "ink_subtle", fs)
        if sl.get("source"):
            cb = self._box(slide, self.W / 3, y - 0.1, self.W / 3, 0.3)
            self._text(cb.text_frame, sl["source"], "ink_subtle", fs, align=PP_ALIGN.CENTER)
        rb = self._box(slide, self.W - self.M - 1.5, y - 0.1, 1.5, 0.3)
        self._text(rb.text_frame, f"{index} / {total}", "ink_subtle", fs, align=PP_ALIGN.RIGHT)

    # ── 레이아웃별 본문 ──
    def _body_bullets(self, slide, blocks, top, bottom, left=None, width=None):
        left = self.M if left is None else left
        width = (self.W - 2 * self.M) if width is None else width
        tb = self._box(slide, left, top, width, bottom - top)
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for b in blocks:
            if b["type"] == "list":
                for it in b["items"]:
                    self._bullet(
                        tf, "".join(r["text"] for r in it["runs"]), it.get("level", 0), first
                    )
                    first = False
                    for ch in it.get("children", []):
                        self._bullet(
                            tf, "".join(r["text"] for r in ch["runs"]), ch.get("level", 1), False
                        )
            elif b["type"] == "paragraph":
                self._bullet(tf, "".join(r["text"] for r in b["runs"]), 0, first, bullet=False)
                first = False
            elif b["type"] == "heading":
                self._bullet(tf, b["text"], 0, first, bold=True, bullet=False)
                first = False

    def _bullet(self, tf, text, level, first, bold=False, bullet=True):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.level = min(level, 4)
        prefix = "· " if bullet else ""
        r = p.add_run()
        r.text = prefix + text
        self._font(r, "ink", self.theme["pptx"].get("body_pt", 14), bold=bold)
        p.space_after = Pt(6)

    def _slide_table(self, slide, block, top, bottom):
        header, rows = block["header"], block["rows"]
        ncol = len(header) or (len(rows[0]) if rows else 1)
        nrow = len(rows) + 1
        tbl_shape = slide.shapes.add_table(
            nrow,
            ncol,
            Inches(self.M),
            Inches(top),
            Inches(self.W - 2 * self.M),
            Inches(min(bottom - top, 0.4 * nrow)),
        )
        table = tbl_shape.table
        _strip_table_style(table)
        tp = self.theme["pptx"]
        for j in range(ncol):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(self.theme, "primary")
            self._cell(
                cell,
                "".join(r["text"] for r in header[j]) if j < len(header) else "",
                "white",
                bold=True,
                size=tp.get("body_pt", 14) - 2,
            )
        for i, row in enumerate(rows, start=1):
            for j in range(ncol):
                cell = table.cell(i, j)
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(self.theme, "surface")
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _rgb(self.theme, "white")
                txt = "".join(r["text"] for r in row[j]) if j < len(row) else ""
                self._cell(cell, txt, "ink", size=tp.get("body_pt", 14) - 3)

    def _cell(self, cell, text, role, bold=False, size=12):
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        self._font(r, role, size, bold=bold)

    def _slide_chart(self, slide, block, top, bottom):
        spec = block["spec"]
        cats = spec.get("categories", [])
        series = spec.get("series", [])
        if len(series) > 4:
            self.warnings.append("차트 시리즈가 5개 이상입니다. 4개 이하로 줄이세요.")
        if len(cats) > 12:
            self.warnings.append("차트 카테고리가 13개 이상입니다.")
        cd = CategoryChartData()
        cd.categories = [str(c) for c in cats]
        for s in series:
            cd.add_series(
                s.get("name", ""), [float(v) if _num(v) else 0 for v in s.get("values", [])]
            )
        ctype = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "stacked": XL_CHART_TYPE.COLUMN_STACKED,
        }.get(spec.get("type", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
        gf = slide.shapes.add_chart(
            ctype,
            Inches(self.M),
            Inches(top),
            Inches(self.W - 2 * self.M),
            Inches(bottom - top - 0.1),
            cd,
        )
        chart = gf.chart
        chart.has_title = False
        try:
            chart.has_legend = len(series) > 1
            if chart.has_legend:
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
        except Exception:  # noqa: BLE001
            pass
        palette = ["primary", "accent", "action", "ink_subtle"]
        for i, plot_series in enumerate(chart.series):
            try:
                plot_series.format.fill.solid()
                plot_series.format.fill.fore_color.rgb = _rgb(self.theme, palette[i % len(palette)])
            except Exception:  # noqa: BLE001
                pass

    def _slide_diagram(self, slide, block, top, bottom):
        rect = (self.M, top, self.W - 2 * self.M, bottom - top)
        res = pptx_shapes.draw_diagram(slide, block["spec"], self.theme, rect)
        self.warnings += res.get("warnings", [])

    def _slide_image(self, slide, block, top, bottom):
        src = Path(block["src"])
        if not src.is_absolute():
            src = self.base_dir / block["src"]
        if src.exists():
            slide.shapes.add_picture(
                str(src), Inches(self.M), Inches(top), width=Inches(self.W - 2 * self.M)
            )
        else:
            self.warnings.append(f"그림 없음: {block['src']}")

    def _first(self, blocks, kind):
        for b in blocks:
            if b["type"] == kind:
                return b
        return None

    # ── 특수 슬라이드 ──
    def _cover(self, prs, fm):
        slide = self._blank(prs)
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.35), Inches(self.H)
        )
        band.fill.solid()
        band.fill.fore_color.rgb = _rgb(self.theme, "primary")
        band.line.fill.background()
        band.shadow.inherit = False
        tb = self._box(slide, 1.0, self.H / 2 - 1.4, self.W - 2.0, 1.4)
        self._text(
            tb.text_frame,
            fm.get("title", "제목 없음"),
            "primary",
            self.theme["pptx"].get("headline_pt", 22) + 14,
            bold=True,
        )
        if fm.get("subtitle"):
            sb = self._box(slide, 1.0, self.H / 2, self.W - 2.0, 0.8)
            self._text(
                sb.text_frame,
                str(fm["subtitle"]),
                "ink_muted",
                self.theme["pptx"].get("headline_pt", 22),
            )
        meta = "  ·  ".join(str(fm[k]) for k in ("org", "author", "date") if fm.get(k))
        mb = self._box(slide, 1.0, self.H - 1.2, self.W - 2.0, 0.5)
        self._text(mb.text_frame, meta, "ink_subtle", self.theme["pptx"].get("body_pt", 14))
        if fm.get("security"):
            sec = self._box(slide, self.W - 2.5, self.M, 2.0, 0.4)
            self._text(
                sec.text_frame,
                str(fm["security"]),
                "ink_subtle",
                self.theme["pptx"].get("footer_pt", 9),
                align=PP_ALIGN.RIGHT,
            )

    def _agenda(self, prs, slides):
        slide = self._blank(prs)
        hb = self._box(slide, self.M, self.M, self.W - 2 * self.M, 0.8)
        self._text(
            hb.text_frame, "목차", "primary", self.theme["pptx"].get("headline_pt", 22), bold=True
        )
        titles = [s["title"] for s in slides if s.get("title")][:8]
        col_n = 1 if len(titles) <= 5 else 2
        per = (len(titles) + col_n - 1) // col_n
        for c in range(col_n):
            tb = self._box(
                slide,
                self.M + c * (self.W - 2 * self.M) / col_n,
                self.M + 1.2,
                (self.W - 2 * self.M) / col_n - 0.2,
                self.H - 2.5,
            )
            tf = tb.text_frame
            tf.word_wrap = True
            for i, title in enumerate(titles[c * per : (c + 1) * per]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                r = p.add_run()
                r.text = f"{c * per + i + 1}.  "
                self._font(r, "primary", self.theme["pptx"].get("body_pt", 14) + 2, bold=True)
                r2 = p.add_run()
                r2.text = title
                self._font(r2, "ink", self.theme["pptx"].get("body_pt", 14) + 2)
                p.space_after = Pt(10)

    # ── 메인 ──
    def render(self, parsed, out_path, template_pptx=None):
        fm = parsed["frontmatter"]
        self.theme["_footer_text"] = fm.get("footer", fm.get("org", ""))
        if template_pptx:
            prs = _open_template(template_pptx, self.base_dir)
        else:
            prs = Presentation()
            prs.slide_width = Inches(self.W)
            prs.slide_height = Inches(self.H)
        slides = parsed["slides"]
        self._cover(prs, fm)
        if fm.get("agenda"):
            self._agenda(prs, slides)
        total = len(slides) + 1 + (1 if fm.get("agenda") else 0)
        base_idx = 2 + (1 if fm.get("agenda") else 0)
        out_meta = []
        for i, sl in enumerate(slides):
            slide = self._blank(prs)
            top, bottom = self._frame(slide, sl, base_idx + i, total)
            self._render_layout(slide, sl, top, bottom)
            if sl.get("notes"):
                slide.notes_slide.notes_text_frame.text = sl["notes"]
            out_meta.append({"index": base_idx + i, "title": sl["title"], "layout": sl["layout"]})
        Path(out_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(out_path)
        return {
            "out_path": str(Path(out_path).resolve()),
            "slide_count": len(prs.slides),
            "slides": out_meta,
            "warnings": self.warnings,
        }

    def _render_layout(self, slide, sl, top, bottom):
        layout = sl["layout"]
        blocks = sl["blocks"]
        if sl.get("columns"):
            cols = sl["columns"]
            n = len(cols)
            gap = 0.3
            cw = (self.W - 2 * self.M - gap * (n - 1)) / n
            for i, col in enumerate(cols):
                left = self.M + i * (cw + gap)
                self._body_bullets(slide, col, top, bottom, left=left, width=cw)
            return
        if layout == "diagram":
            blk = self._first(blocks, "diagram")
            if blk:
                self._slide_diagram(slide, blk, top, bottom)
            return
        if layout == "chart":
            blk = self._first(blocks, "chart")
            if blk:
                self._slide_chart(slide, blk, top, bottom)
            return
        if layout in ("table", "timeline"):
            blk = self._first(blocks, "table")
            if not blk and layout == "timeline":
                tl = self._first(blocks, "timeline")
                if tl:
                    blk = _timeline_table(tl)
            if blk:
                self._slide_table(slide, blk, top, bottom)
            return
        if layout == "image":
            blk = self._first(blocks, "image")
            if blk:
                self._slide_image(slide, blk, top, bottom)
            return
        # section / message / closing → 글머리표
        self._body_bullets(slide, blocks, top, bottom)


def _num(v):
    try:
        float(v)
        return True
    except (TypeError, ValueError):
        return False


def _timeline_table(tl):
    tasks = tl["spec"].get("tasks", [])
    r = lambda s: [{"text": str(s), "bold": False, "italic": False, "code": False}]  # noqa: E731
    return {
        "type": "table",
        "header": [r("항목"), r("시작"), r("종료"), r("비고")],
        "rows": [
            [
                r(t.get("label", "")),
                r(t.get("start", "")),
                r(t.get("end", "")),
                r(t.get("group", "")),
            ]
            for t in tasks
        ],
        "align": ["left", "left", "left", "left"],
        "caption": None,
    }


def _strip_table_style(table):
    tblPr = table._tbl.tblPr
    for a in ("firstRow", "bandRow"):
        tblPr.set(a, "0")


def _open_template(template_pptx, base_dir):
    p = Path(template_pptx)
    if not p.is_absolute() and base_dir:
        p = Path(base_dir) / template_pptx
    prs = Presentation(str(p))
    # 기존 슬라이드 제거(마스터·레이아웃은 유지). 관계도 끊어 orphan part 가 저장되지 않게 한다.
    xml_slides = prs.slides._sldIdLst
    for sid in list(xml_slides):
        rid = sid.get(qn("r:id"))
        if rid:
            try:
                prs.part.drop_rel(rid)
            except KeyError:
                pass
        xml_slides.remove(sid)
    return prs


def render_pptx(parsed, out_path, theme, base_dir=None, template_pptx=None):
    return DeckRenderer(theme, base_dir).render(parsed, out_path, template_pptx)
