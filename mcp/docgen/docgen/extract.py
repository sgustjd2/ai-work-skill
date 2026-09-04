"""기존 docx·pptx → Markdown, 회사 템플릿 pptx → 테마 JSON, 레이아웃 목록.

회사 양식의 관용 구조를 스킬이 배우고, 공식 색·폰트를 테마로 뽑는 용도.
"""

from __future__ import annotations

import json
import zipfile
from pathlib import Path

from docx import Document
from docx.oxml.ns import qn as docx_qn
from pptx import Presentation


def _para_md(p) -> str:
    style = (p.style.name or "").lower() if p.style else ""
    text = p.text.strip()
    if not text:
        return ""
    if style.startswith("heading"):
        try:
            level = int(style.split()[-1])
        except ValueError:
            level = 2
        return "#" * min(level, 6) + " " + text
    if "list" in style:
        return "- " + text
    return text


def _table_md(tbl) -> list[str]:
    rows = tbl.rows
    if not rows:
        return []
    out = []
    header = [c.text.strip().replace("\n", " ") for c in rows[0].cells]
    out.append("| " + " | ".join(header) + " |")
    out.append("|" + "|".join(["---"] * len(header)) + "|")
    for r in rows[1:]:
        out.append("| " + " | ".join(c.text.strip().replace("\n", " ") for c in r.cells) + " |")
    return out


def docx_to_md(path: str) -> dict:
    doc = Document(path)
    lines: list[str] = []
    headings, tables, images = [], 0, 0
    body = doc.element.body
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    for child in body.iterchildren():
        if child.tag == docx_qn("w:p"):
            p = Paragraph(child, doc)
            md = _para_md(p)
            if md:
                lines.append(md)
                if md.startswith("#"):
                    headings.append(md.lstrip("# ").strip())
            if "w:drawing" in child.xml or "pic:pic" in child.xml:
                images += 1
        elif child.tag == docx_qn("w:tbl"):
            t = Table(child, doc)
            lines.append("")
            lines += _table_md(t)
            lines.append("")
            tables += 1
    md = "\n".join(lines).strip() + "\n"
    return {
        "markdown": md,
        "structure": {"headings": headings, "tables": tables, "images": images},
        "warnings": [],
    }


def pptx_to_md(path: str) -> dict:
    prs = Presentation(path)
    lines: list[str] = []
    tables = 0
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()
        lines.append(f"# {title or f'슬라이드 {i}'}")
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_table:
                tbl = shape.table
                header = [c.text.strip() for c in tbl.rows[0].cells] if tbl.rows else []
                if header:
                    lines.append("| " + " | ".join(header) + " |")
                    lines.append("|" + "|".join(["---"] * len(header)) + "|")
                    for r in list(tbl.rows)[1:]:
                        lines.append("| " + " | ".join(c.text.strip() for c in r.cells) + " |")
                    tables += 1
            elif shape.has_text_frame and shape.text.strip():
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t and t != title:
                        lines.append(f"- {t}")
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                lines.append(f"<!-- note: {note} -->")
        lines.append("\n---")
    md = "\n".join(lines).strip("-\n ") + "\n"
    return {
        "markdown": md,
        "structure": {"slides": len(prs.slides), "tables": tables},
        "warnings": [],
    }


# ─── 테마 추출 ───────────────────────────────────────────────────────────────
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _clr(el) -> str | None:
    if el is None:
        return None
    srgb = el.find(f"{{{_A}}}srgbClr")
    if srgb is not None:
        return "#" + srgb.get("val").upper()
    sysc = el.find(f"{{{_A}}}sysClr")
    if sysc is not None and sysc.get("lastClr"):
        return "#" + sysc.get("lastClr").upper()
    return None


def theme_from_pptx(path: str, name: str, out_dir: str = "themes") -> dict:
    import xml.etree.ElementTree as ET

    warnings: list[str] = []
    with zipfile.ZipFile(path) as z:
        theme_xml = next((n for n in z.namelist() if n.startswith("ppt/theme/theme")), None)
        if not theme_xml:
            raise ValueError("theme1.xml 을 찾을 수 없습니다.")
        root = ET.fromstring(z.read(theme_xml))

    ns = f"{{{_A}}}"
    scheme = root.find(f"{ns}themeElements/{ns}clrScheme")
    colors_raw = {}
    for child in list(scheme) if scheme is not None else []:
        key = child.tag.split("}")[-1]
        colors_raw[key] = _clr(child)
    font_scheme = root.find(f"{ns}themeElements/{ns}fontScheme")
    minor = font_scheme.find(f"{ns}minorFont/{ns}latin") if font_scheme is not None else None
    major_ea = font_scheme.find(f"{ns}majorFont/{ns}ea") if font_scheme is not None else None

    # 매핑(초안): accent1→primary, accent2→accent, dk1→ink, lt2→surface, hlink→action
    default = json.loads((Path(_repo_themes()) / "datasolution.json").read_text(encoding="utf-8"))
    colors = dict(default["colors"])
    mapping = {
        "primary": "accent1",
        "accent": "accent2",
        "action": "hlink",
        "ink": "dk1",
        "surface": "lt2",
    }
    for role, src in mapping.items():
        v = colors_raw.get(src)
        if v:
            colors[role] = v
        else:
            warnings.append(f"{role} ← {src} 없음, 기본값 유지")
    warnings.append("색 매핑은 초안입니다. 확인 필요(accent1→primary 등).")

    fonts = dict(default["fonts"])
    if minor is not None and minor.get("typeface"):
        fonts["office_latin"] = minor.get("typeface")
    if major_ea is not None and major_ea.get("typeface"):
        fonts["office_ko"] = major_ea.get("typeface")

    prs = Presentation(path)
    theme = {
        "name": name,
        "source": f"{Path(path).name} theme1.xml 추출",
        "colors": colors,
        "fonts": fonts,
        "docx": default["docx"],
        "pptx": {
            **default["pptx"],
            "width_in": round(prs.slide_width / 914400, 3),
            "height_in": round(prs.slide_height / 914400, 3),
            "layout_map": {},
        },
    }
    out = Path(out_dir) / f"{name}.json"
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(theme, ensure_ascii=False, indent=2), encoding="utf-8")
    return {
        "out_path": str(out.resolve()),
        "colors": {k: colors[k] for k in mapping},
        "warnings": warnings,
    }


def layouts(path: str) -> dict:
    prs = Presentation(path)
    out = []
    for layout in prs.slide_layouts:
        placeholders = [
            {
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
                "name": ph.name,
            }
            for ph in layout.placeholders
        ]
        out.append({"name": layout.name, "placeholders": placeholders})
    return {
        "layouts": out,
        "size": {
            "width_in": round(prs.slide_width / 914400, 3),
            "height_in": round(prs.slide_height / 914400, 3),
        },
        "warnings": [],
    }


def _repo_themes() -> str:
    env = __import__("os").environ.get("AI_WORK_SKILL_ROOT")
    if env:
        return str(Path(env) / "themes")
    return str(Path(__file__).resolve().parents[3] / "themes")
